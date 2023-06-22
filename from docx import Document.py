from docx import Document
from pptx import Presentation

def convert_docx_to_pptx(docx_file, pptx_file):
    # Open the document file
    doc = Document(docx_file)

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Iterate through each paragraph in the document
    for para in doc.paragraphs:
        # Add a new slide to the presentation
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Get the title and content placeholders from the slide
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        # Set the paragraph text as the slide title
        title_placeholder.text = para.text

        # Create a new paragraph for the slide content
        p = content_placeholder.text_frame.add_paragraph()

        # Add the paragraph text as the slide content
        p.text = para.text

    # Save the PowerPoint presentation
    prs.save(pptx_file)

# Provide the file paths for the document and PowerPoint files
docx_file_path = r"C:\Users\prath\DATABASE MANAGEMENT SYSTEM FINAL\Title.docx"
pptx_file_path = r"C:\Users\prath\DATABASE MANAGEMENT SYSTEM FINAL\DATABASE MANAGEMENT SYSTEM PPT.pptx"

# Convert the document content into PowerPoint slides
convert_docx_to_pptx(docx_file_path, pptx_file_path)
